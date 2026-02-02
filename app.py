import streamlit as st
import pandas as pd
import numpy as np
import io
import validator
import kpi_engine
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from datetime import datetime

# Page Config
st.set_page_config(page_title="Supply Chain Risk Cockpit", layout="wide", page_icon="✈️")

# ---------------------------------------------------------
# Styles
# ---------------------------------------------------------
st.markdown("""
<style>
    .metric-card {
        background-color: #f0f2f6;
        padding: 15px;
        border-radius: 10px;
        text-align: center;
    }
    .main-header {
        text-align: center;
    }
    .social-link {
        text-decoration: none;
        color: #0072b1;
        font-weight: bold;
        margin-right: 15px;
    }
</style>
""", unsafe_allow_html=True)

# ---------------------------------------------------------
# Helper Functions
# ---------------------------------------------------------
@st.cache_data
def load_excel(file_bytes):
    xls = pd.ExcelFile(io.BytesIO(file_bytes))
    dfs = {}
    for sh in xls.sheet_names:
        dfs[sh] = pd.read_excel(xls, sheet_name=sh)
    return dfs

def enrich_master(df, dfs):
    """Joins master data (unit_revenue, cogs, etc.)"""
    if "Master_Data" not in dfs:
        df["unit_revenue"] = 1.0
        df["unit_cogs"] = 0.5
        return df
    
    m = dfs["Master_Data"].copy()
    if "location" in m.columns:
        m = m.drop_duplicates(subset=["sku", "location"])
        df = df.merge(m[["sku", "location", "unit_revenue", "unit_cogs"]], on=["sku", "location"], how="left")
    else:
        m = m.drop_duplicates(subset=["sku"])
        df = df.merge(m[["sku", "unit_revenue", "unit_cogs"]], on=["sku"], how="left")
        
    df["unit_revenue"] = df["unit_revenue"].fillna(1.0)
    df["unit_cogs"] = df["unit_cogs"].fillna(0.5)
    return df

def generate_ppt(summary, tot_rar, skus_stockout, avg_fill, safety_breaches):
    """Generate PowerPoint presentation with dashboard results"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Title Slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    title.text = "Supply Chain Risk Dashboard"
    subtitle.text = f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M')}"
    
    # KPI Summary Slide
    kpi_slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = kpi_slide.shapes.title
    title.text = "Executive KPI Summary"
    
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(0.8)
    
    metrics = [
        ("💰 Revenue at Risk", f"${tot_rar:,.0f}"),
        ("⚠️ SKUs with Stockouts", f"{int(skus_stockout)}"),
        ("📉 Average Fill Rate", f"{avg_fill*100:.1f}%"),
        ("🛡️ Safety Stock Breaches", f"{int(safety_breaches)}")
    ]
    
    for i, (label, value) in enumerate(metrics):
        textbox = kpi_slide.shapes.add_textbox(left, top + i*height, width, height)
        text_frame = textbox.text_frame
        text_frame.text = f"{label}: {value}"
        text_frame.paragraphs[0].font.size = Pt(18)
        text_frame.paragraphs[0].font.bold = True
    
    # Top Risks Slide
    risks_slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = risks_slide.shapes.title
    title.text = "Top 10 At-Risk SKU-Locations"
    
    top_risks = summary.nlargest(10, 'revenue_at_risk')[['sku', 'location', 'revenue_at_risk', 'fill_rate']]
    
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(4)
    
    table = risks_slide.shapes.add_table(len(top_risks) + 1, 4, left, top, width, height).table
    
    # Header
    headers = ['SKU', 'Location', 'Revenue at Risk', 'Fill Rate']
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
    
    # Data
    for row_idx, (_, row) in enumerate(top_risks.iterrows(), start=1):
        table.cell(row_idx, 0).text = str(row['sku'])
        table.cell(row_idx, 1).text = str(row['location'])
        table.cell(row_idx, 2).text = f"${row['revenue_at_risk']:,.0f}"
        table.cell(row_idx, 3).text = f"{row['fill_rate']*100:.1f}%"
    
    # Save to BytesIO
    ppt_io = io.BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)
    return ppt_io

# ---------------------------------------------------------
# Views
# ---------------------------------------------------------

def show_landing_page():
    # Header Image
    if os.path.exists("assets/header.jpg"):
        st.image("assets/header.jpg", use_container_width=True)
    
    st.title("Supply Chain Decision Tool")
    
    col1, col2 = st.columns([2, 1])
    
    with col1:
        st.subheader("How to Use")
        st.markdown("""
        1. **Download the Template**: Get the required Excel format below.
        2. **Fill Your Data**: Enter Inventory, Demand, and Supply data.
        3. **Upload**: Go to the 'Tool' page and upload your file.
        4. **Analyze**: View Risk calculations and actionable recommendations.
        """)
        
        # Download Button
        template_path = "control_tower_input_with_help.xlsx"
        if os.path.exists(template_path):
            with open(template_path, "rb") as f:
                st.download_button(
                    label="📥 Download Excel Template",
                    data=f,
                    file_name="control_tower_input_with_help.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                )
        else:
            st.warning("Template file not found. Please run make_template.py")

        st.write("") # Spacer
        if st.button("🚀 Launch Tool", type="primary"):
            st.session_state.page = "Tool"
            st.rerun()

    with col2:
        st.subheader("Connect with Me")
        st.markdown("""
        Built by **Ankit Nanda**.
        
        - [Substack: Supply Chain Insider](https://substack.com/@ankitnanda)
        - [LinkedIn Profile](https://www.linkedin.com/in/ankitnanda/)
        """)

def show_about():
    st.title("About This Tool")
    
    st.markdown("""
    ## 🎯 What This Tool Does
    
    The **Supply Chain Decision Tool** helps you:
    
    - **Calculate Key Metrics**: Net Available Inventory (NAI), Projected On-Hand (POH), Fill Rate, Revenue at Risk
    - **Identify Risks**: Stockout flags, safety stock breaches, first-week-of-risk alerts
    - **Run Scenarios**: Test demand uplifts (0-50%) and supply delays (0-8 weeks)
    - **Get Actionable Insights**: Prioritized recommendations (Expedite, Replenish, Promo)
    - **Export Results**: Download dashboard as PowerPoint for presentations
    
    ### How It Works
    
    1. **Upload Your Data**: Use our Excel template with Inventory, Demand, and Supply sheets
    2. **Analyze**: The tool calculates KPIs using deterministic logic (no ML black boxes)
    3. **Scenario Test**: Adjust parameters to see "what-if" impacts
    4. **Act**: Get prioritized recommendations based on risk severity
    
    ## 🔒 Privacy & Data Security
    
    **Your data is 100% private:**
    
    - ✅ All processing happens **client-side** (in your browser)
    - ✅ **No data is stored** on any server
    - ✅ **No data is transmitted** to external services
    - ✅ **No authentication required** - no user accounts or tracking
    - ✅ **Open source** - you can verify the code yourself
    
    Upload your confidential supply chain data with complete confidence.
    
    ## 💡 Use Cases
    
    ### Pharma Supply Chain
    - Track critical medications (Antibiotics, Analgesics, etc.)
    - Prevent stockouts that impact patient care
    - Optimize inventory to reduce holding costs
    
    ### Multi-Location Distribution
    - Manage inventory across DCs and warehouses
    - Identify imbalances and redistribution opportunities
    - Coordinate supply across network
    
    ### S&OP Planning
    - Weekly bucket analysis aligned with planning cycles
    - Scenario modeling for demand/supply variability
    - Executive dashboards for leadership reviews
    
    ## 🛠️ Built With
    
    - **Python 3.13+**
    - **Streamlit** - Web framework
    - **Pandas** - Data processing
    - **OpenPyXL** - Excel handling
    - **python-pptx** - PowerPoint export
    
    ## 📞 Contact
    
    Built by **Ankit Nanda** - Supply Chain Professional
    
    - [Substack: Supply Chain Insider](https://substack.com/@ankitnanda)
    - [LinkedIn](https://www.linkedin.com/in/ankitnanda/)
    - [GitHub Repository](https://github.com/akn3107/scm-decision-cockpit)
    
    ---
    
    **Open Source | MIT License | Built with ❤️ for supply chain professionals**
    """)

def show_cockpit():
    st.header("Risk Control Tower")
    
    # Sidebar scenarios
    demand_uplift = st.sidebar.slider("Demand Uplift (%)", 0, 50, 0, 5) / 100.0
    supply_delay = st.sidebar.slider("Supply Delay (Weeks)", 0, 8, 0, 1)
    horizon = st.sidebar.slider("Planning Horizon (Weeks)", 4, 16, 8, 1)

    uploaded_file = st.file_uploader("Upload 'control_tower_input_with_help.xlsx'", type=["xlsx"])
    
    if not uploaded_file:
        st.info("👆 Please upload the data file to proceed.")
        return

    # Data Processing
    try:
        file_bytes = uploaded_file.read()
        dfs = load_excel(file_bytes)
        
        ok, errors = validator.validate_data(dfs)
        if not ok:
            st.error("❌ Data Validation Failed")
            for e in errors:
                st.write(f"- {e}")
            st.stop()
            
        inv = dfs["Inventory"]
        demand = dfs["Demand_Plan"]
        supply = dfs["Supply_Plan"]
        
        with st.spinner("Computing Scenarios..."):
            summary, detail = kpi_engine.compute_kpis(
                inv, demand, supply, 
                horizon_weeks=horizon,
                demand_uplift_pct=demand_uplift,
                supply_delay_weeks=supply_delay
            )
            
        summary = enrich_master(summary, dfs)
        summary["revenue_at_risk"] = summary["total_unmet"] * summary["unit_revenue"]
        
        # Executive Metrics
        tot_rar = summary["revenue_at_risk"].sum()
        skus_stockout = summary["stockout_flag"].sum()
        avg_fill = summary["fill_rate"].mean()
        safety_breaches = summary["safety_breach_flag"].sum()
        
        c1, c2, c3, c4 = st.columns(4)
        c1.metric("💰 Revenue at Risk", f"${tot_rar:,.0f}")
        c2.metric("⚠️ SKUs with Stockouts", int(skus_stockout))
        c3.metric("📉 Avg Fill Rate", f"{avg_fill*100:.1f}%")
        c4.metric("🛡️ Safety Breaches", int(safety_breaches))
        
        # PPT Export Button
        st.write("")  # Spacer
        ppt_data = generate_ppt(summary, tot_rar, skus_stockout, avg_fill, safety_breaches)
        st.download_button(
            label="📊 Download as PowerPoint",
            data=ppt_data,
            file_name=f"supply_chain_dashboard_{datetime.now().strftime('%Y%m%d_%H%M')}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
        
        st.divider()
        
        # Actions
        st.subheader("🔥 Top Actions")
        def recommend(row):
            if row["stockout_flag"]: return "🟥 EXPEDITE"
            if row["safety_breach_flag"]: return "🟧 REPLENISH"
            if row["min_poh"] > row["safety_stock_qty"]*3 and row["min_poh"] > 10000: return "🟦 PROMO"
            return "✅ OK"
            
        actions = summary.copy()
        actions["Recommendation"] = actions.apply(recommend, axis=1)
        actions = actions.sort_values(["revenue_at_risk", "first_stockout_week"], ascending=[False, True])
        
        st.dataframe(
            actions[[
                "sku", "location", "Recommendation", 
                "revenue_at_risk", "fill_rate", 
                "first_stockout_week", "first_safety_breach_week"
            ]].style.format({
                "revenue_at_risk": "${:,.0f}",
                "fill_rate": "{:.1%}"
            }).applymap(lambda x: "color: red; font-weight: bold" if "EXPEDITE" in str(x) else ("color: orange; font-weight: bold" if "REPLENISH" in str(x) else ""), subset=["Recommendation"]),
            use_container_width=True
        )
        
        # Drilldown
        st.divider()
        st.subheader("🔍 Drilldown")
        sku = st.selectbox("Select SKU", summary["sku"].unique())
        loc = st.selectbox("Select Location", summary[summary["sku"]==sku]["location"].unique())
        
        drill = detail[(detail["sku"]==sku) & (detail["location"]==loc)].copy()
        st.line_chart(drill.set_index("week_start")[["on_hand_qty", "forecast_qty", "supply_qty"]])
        st.dataframe(drill[["week_start", "forecast_qty", "supply_qty", "NAI", "POH", "served_qty", "unmet_qty"]].style.format("{:,.0f}"), use_container_width=True)

    except Exception as e:
        st.error(f"Error: {str(e)}")

# ---------------------------------------------------------
# Main Navigation
# ---------------------------------------------------------
def main():
    if "page" not in st.session_state:
        st.session_state.page = "Home"

    st.sidebar.title("Navigation")
    
    # Callback to sync sidebar with session state
    def on_change():
        st.session_state.page = st.session_state.nav_radio

    selection = st.sidebar.radio(
        "Go to", 
        ["Home", "Tool", "About"], 
        key="nav_radio",
        index=0 if st.session_state.page == "Home" else (1 if st.session_state.page == "Tool" else 2),
        on_change=on_change
    )

    if st.session_state.page == "Home":
        show_landing_page()
    elif st.session_state.page == "Tool":
        show_cockpit()
    elif st.session_state.page == "About":
        show_about()

if __name__ == "__main__":
    main()
